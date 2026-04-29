from __future__ import annotations

from io import BytesIO
from pathlib import Path
from uuid import UUID

import pdfplumber
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

from rag_hybrid.app_logging import get_logger
from rag_hybrid.assets import save_binary_asset
from rag_hybrid.config import get_settings
from rag_hybrid.models import ExtractedChunk, ExtractedTable

logger = get_logger()
CHUNK_SIZE = get_settings().extraction.chunk_size


def _normalize_cell(cell: object) -> str:
    return "" if cell is None else str(cell).strip()


def _table_summary(rows: list[list[str]]) -> str:
    preview = " | ".join(rows[0]) if rows else "No headers"
    return f"Table extracted with {len(rows)} rows. Header preview: {preview}"


def _table_as_text(rows: list[list[str]], max_rows: int | None = None) -> str:
    if not rows:
        return ""

    selected_rows = rows if max_rows is None else rows[:max_rows]
    rendered_rows = [" | ".join(cell for cell in row if cell is not None).strip() for row in selected_rows]
    return "\n".join(row for row in rendered_rows if row)


def _append_table_chunk(
    chunks: list[ExtractedChunk],
    rows: list[list[str]],
    section: str,
    source: str,
    page_number: int | None = None,
) -> None:
    summary = _table_summary(rows)
    table_text = _table_as_text(rows)
    chunk_content = "\n".join(
        [
            f"Structured table from {section}",
            summary,
            table_text,
        ]
    ).strip()

    chunks.append(
        ExtractedChunk(
            content=chunk_content,
            content_type="table",
            section=section,
            chunk_id=len(chunks) + 1,
            source=source,
            page_number=page_number,
            section_title=section,
        )
    )


def _append_image_chunk(
    chunks: list[ExtractedChunk],
    document_id: UUID,
    source: str,
    section: str,
    image_bytes: bytes,
    extension: str,
    description: str,
    page_number: int | None = None,
) -> None:
    asset_path = save_binary_asset(document_id, source, extension, image_bytes)
    chunks.append(
        ExtractedChunk(
            content=description,
            content_type="image",
            section=section,
            chunk_id=len(chunks) + 1,
            source=source,
            page_number=page_number,
            section_title=section,
            metadata={
                "asset_path": asset_path,
                "asset_type": "image",
                "extension": extension,
            },
        )
    )


def _chunk_text(text: str, source: str, section: str, start_index: int, size: int = CHUNK_SIZE, page_number: int | None = None) -> list[ExtractedChunk]:
    normalized = " ".join(text.split())
    if not normalized:
        return []

    chunks = []
    cursor = 0
    chunk_id = start_index
    while cursor < len(normalized):
        piece = normalized[cursor:cursor + size].strip()
        if piece:
            chunks.append(
                ExtractedChunk(
                    content=piece,
                    content_type="text",
                    section=section,
                    chunk_id=chunk_id,
                    source=source,
                    page_number=page_number,
                    section_title=section,
                )
            )
            chunk_id += 1
        cursor += size
    return chunks


def extract_file(file_name: str, file_bytes: bytes, document_id: UUID) -> tuple[list[ExtractedChunk], list[ExtractedTable]]:
    suffix = Path(file_name).suffix.lower()
    logger.info("Dispatching extraction for file '%s' with extension '%s'.", file_name, suffix)

    if suffix == ".pdf":
        return extract_pdf(file_name, file_bytes, document_id)
    if suffix == ".docx":
        return extract_docx(file_name, file_bytes, document_id)
    if suffix == ".pptx":
        return extract_pptx(file_name, file_bytes, document_id)

    raise ValueError(f"Unsupported file type: {suffix}")


def extract_pdf(file_name: str, file_bytes: bytes, document_id: UUID) -> tuple[list[ExtractedChunk], list[ExtractedTable]]:
    chunks: list[ExtractedChunk] = []
    tables: list[ExtractedTable] = []
    logger.info("Starting PDF extraction for '%s'.", file_name)

    reader = PdfReader(BytesIO(file_bytes))
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        chunks.extend(_chunk_text(text, file_name, f"Page {page_number}", len(chunks) + 1, page_number=page_number))

    with pdfplumber.open(BytesIO(file_bytes)) as pdf:
        for page_number, page in enumerate(pdf.pages, start=1):
            for table_number, raw_table in enumerate(page.extract_tables() or [], start=1):
                rows = [[_normalize_cell(cell) for cell in row] for row in raw_table if row]
                if not rows:
                    continue

                summary = _table_summary(rows)
                section = f"Page {page_number} Table {table_number}"
                _append_table_chunk(chunks, rows, section, file_name, page_number)
                tables.append(
                    ExtractedTable(
                        table_id=f"{document_id}-pdf-{page_number}-{table_number}",
                        document_id=document_id,
                        source=file_name,
                        section=section,
                        rows=rows,
                        summary=summary,
                        headers=rows[0] if rows else [],
                        page_number=page_number,
                        section_title=section,
                        table_name=section,
                    )
                )

    for page_number, page in enumerate(reader.pages, start=1):
        for image_number, image_file in enumerate(getattr(page, "images", []), start=1):
            image_bytes = getattr(image_file, "data", None)
            image_name = getattr(image_file, "name", f"page_{page_number}_{image_number}.png")
            extension = Path(image_name).suffix or ".png"
            if image_bytes:
                _append_image_chunk(
                    chunks,
                    document_id,
                    file_name,
                    f"Page {page_number}",
                    image_bytes,
                    extension,
                    f"Image extracted from page {page_number}, image {image_number} in {file_name}",
                    page_number,
                )

    logger.info("Completed PDF extraction for '%s': %s chunks, %s tables.", file_name, len(chunks), len(tables))
    return chunks, tables


def extract_docx(file_name: str, file_bytes: bytes, document_id: UUID) -> tuple[list[ExtractedChunk], list[ExtractedTable]]:
    chunks: list[ExtractedChunk] = []
    tables: list[ExtractedTable] = []
    logger.info("Starting DOCX extraction for '%s'.", file_name)

    document = Document(BytesIO(file_bytes))
    body_text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
    chunks.extend(_chunk_text(body_text, file_name, "Document Body", len(chunks) + 1))

    for table_number, table in enumerate(document.tables, start=1):
        rows = [[_normalize_cell(cell.text) for cell in row.cells] for row in table.rows]
        summary = _table_summary(rows)
        section = f"Table {table_number}"
        _append_table_chunk(chunks, rows, section, file_name)
        tables.append(
            ExtractedTable(
                table_id=f"{document_id}-docx-{table_number}",
                document_id=document_id,
                source=file_name,
                section=section,
                rows=rows,
                summary=summary,
                headers=rows[0] if rows else [],
                section_title=section,
                table_name=section,
            )
        )

    image_count = 0
    for relation in document.part.rels.values():
        target_ref = getattr(relation, "target_ref", "")
        if "image" not in target_ref:
            continue

        if getattr(relation, "is_external", False):
            logger.info(
                "Skipping external image relationship '%s' in '%s'.",
                target_ref,
                file_name,
            )
            continue

        try:
            target_part = relation.target_part
        except ValueError:
            logger.info(
                "Skipping non-embedded image relationship '%s' in '%s'.",
                target_ref,
                file_name,
            )
            continue

        if target_part is None:
            continue

        if "image" in target_ref:
            image_count += 1
            image_bytes = getattr(target_part, "blob", None)
            extension = Path(getattr(target_part, "partname", f"image_{image_count}.png")).suffix or ".png"
            if image_bytes:
                _append_image_chunk(
                    chunks,
                    document_id,
                    file_name,
                    "Images",
                    image_bytes,
                    extension,
                    f"Image extracted from Word document {file_name}, image {image_count}",
                )

    logger.info("Completed DOCX extraction for '%s': %s chunks, %s tables.", file_name, len(chunks), len(tables))
    return chunks, tables


def extract_pptx(file_name: str, file_bytes: bytes, document_id: UUID) -> tuple[list[ExtractedChunk], list[ExtractedTable]]:
    chunks: list[ExtractedChunk] = []
    tables: list[ExtractedTable] = []
    logger.info("Starting PPTX extraction for '%s'.", file_name)

    presentation = Presentation(BytesIO(file_bytes))
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text_parts: list[str] = []

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                slide_text_parts.append(shape.text.strip())

            if getattr(shape, "has_table", False):
                rows = [
                    [_normalize_cell(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                summary = _table_summary(rows)
                section = f"Slide {slide_number} Table"
                _append_table_chunk(chunks, rows, section, file_name, slide_number)
                tables.append(
                    ExtractedTable(
                        table_id=f"{document_id}-pptx-{slide_number}-{len(tables) + 1}",
                        document_id=document_id,
                        source=file_name,
                        section=section,
                        rows=rows,
                        summary=summary,
                        headers=rows[0] if rows else [],
                        page_number=slide_number,
                        section_title=section,
                        table_name=section,
                    )
                )

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = getattr(shape, "image", None)
                if image is not None:
                    extension = f".{image.ext}" if getattr(image, "ext", None) else ".png"
                    _append_image_chunk(
                        chunks,
                        document_id,
                        file_name,
                        f"Slide {slide_number}",
                        image.blob,
                        extension,
                        f"Image extracted from slide {slide_number} in presentation {file_name}",
                        slide_number,
                    )

        slide_text = "\n".join(slide_text_parts)
        chunks.extend(_chunk_text(slide_text, file_name, f"Slide {slide_number}", len(chunks) + 1, page_number=slide_number))

    logger.info("Completed PPTX extraction for '%s': %s chunks, %s tables.", file_name, len(chunks), len(tables))
    return chunks, tables
